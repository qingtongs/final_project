"""
Generate Chinese-only PPT presentation for the Animal Detection project.
"""
from __future__ import annotations

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

PROJECT_ROOT = Path(__file__).resolve().parents[1]
OUTPUT_DIR = PROJECT_ROOT / "outputs" / "final_deliverables"
PPT_PATH = OUTPUT_DIR / "animal_detection_project_presentation_zh.pptx"

# Colors
DARK_BG = RGBColor(0x0E, 0x1A, 0x2B)
BLUE = RGBColor(0x24, 0x5A, 0x9A)
CYAN = RGBColor(0x32, 0xA6, 0xB8)
GREEN = RGBColor(0x4B, 0x9B, 0x68)
ORANGE = RGBColor(0xE2, 0x8A, 0x32)
RED = RGBColor(0xC9, 0x4F, 0x4F)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x17, 0x20, 0x33)
MUTED = RGBColor(0x5E, 0x6A, 0x7D)
PAPER = RGBColor(0xF7, 0xF9, 0xFC)
LINE = RGBColor(0xD8, 0xE2, 0xEF)
ZH_FONT = "Microsoft YaHei"


def add_bg(slide, color=PAPER):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    return bg


def add_text(slide, text, left, top, width, height, size=18, color=INK,
             bold=False, align=PP_ALIGN.LEFT, font=ZH_FONT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.alignment = align
    return txBox


def add_multiline(slide, lines, left, top, width, size=18, color=INK):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(6))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = ZH_FONT
        p.space_after = Pt(4)
    return txBox


def add_bullets(slide, items, left, top, width, size=15, color=INK):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(5.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"▸ {item}"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = ZH_FONT
        p.space_after = Pt(6)
    return txBox


def add_rect(slide, left, top, width, height, fill, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_rounded(slide, left, top, width, height, fill, line=LINE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(1)
    return shape


def add_footer(slide, n, total=10):
    add_rect(slide, 0.5, 7.05, 12.3, 0.01, LINE)
    add_text(slide, "YOLO11s 动物检测与计数 | 深度学习期末项目",
             0.5, 7.1, 9.5, 0.35, size=9, color=MUTED)
    add_text(slide, f"{n}/{total}", 11.5, 7.1, 0.8, 0.35, size=11, color=BLUE,
             bold=True, align=PP_ALIGN.RIGHT)


def add_kicker(slide, text):
    marker = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    Inches(0.5), Inches(0.45), Inches(0.15), Inches(0.15))
    marker.fill.solid()
    marker.fill.fore_color.rgb = CYAN
    marker.line.fill.background()
    add_text(slide, text, 0.8, 0.35, 6.0, 0.35, size=11, color=BLUE, bold=True)


def add_title_block(slide, title, subtitle=""):
    add_text(slide, title, 0.5, 0.85, 12.0, 0.8, size=30, color=INK, bold=True)
    if subtitle:
        add_text(slide, subtitle, 0.55, 1.55, 11.0, 0.4, size=15, color=MUTED)


def add_metric(slide, value, label, left, top, color=BLUE):
    add_rounded(slide, left, top, 2.2, 1.5, WHITE, LINE)
    add_text(slide, value, left + 0.2, top + 0.15, 1.8, 0.45, size=26, color=color, bold=True)
    add_text(slide, label, left + 0.2, top + 0.65, 1.8, 0.5, size=12, color=MUTED)


def add_step(slide, num, title, desc, left, top, color=BLUE):
    add_rounded(slide, left, top, 2.0, 1.35, WHITE, LINE)
    add_rect(slide, left, top, 2.0, 0.08, color)
    add_text(slide, f"步骤 {num}", left + 0.15, top + 0.18, 1.7, 0.25, size=10, color=color, bold=True)
    add_text(slide, title, left + 0.15, top + 0.45, 1.7, 0.28, size=15, color=INK, bold=True)
    add_text(slide, desc, left + 0.15, top + 0.78, 1.7, 0.45, size=9, color=MUTED)


def add_bar(slide, label, value, max_val, left, top, color=GREEN):
    bw = 5.0
    add_text(slide, label, left, top, 2.2, 0.28, size=12, color=INK)
    add_rect(slide, left + 2.4, top + 0.06, bw, 0.16, RGBColor(0xE8, 0xEE, 0xF6))
    if value > 0:
        add_rect(slide, left + 2.4, top + 0.06, bw * value / max_val, 0.16, color)
    add_text(slide, f"{value:.1f}" if isinstance(value, float) else str(value),
             left + 2.4 + bw + 0.15, top, 0.6, 0.28, size=12, color=color, bold=True)


# ============================================================
# SLIDE 1: 封面
# ============================================================
def slide01(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, PAPER)
    add_rect(slide, 0, 0, 4.3, 7.5, DARK_BG)

    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(3.3), Inches(1.8))
    tf = txBox.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = "动物检测与计数系统"
    p.font.size = Pt(36); p.font.color.rgb = WHITE; p.font.bold = True; p.font.name = ZH_FONT
    p2 = tf.add_paragraph(); p2.text = "YOLO11s + 自定义20类微调模型"
    p2.font.size = Pt(16); p2.font.color.rgb = RGBColor(0xBB, 0xD2, 0xF0); p2.font.name = ZH_FONT

    add_metric(slide, "20", "动物类别", 5.3, 1.5, BLUE)
    add_metric(slide, "0.978", "最终 mAP50", 7.8, 1.5, GREEN)
    add_metric(slide, "90.94", "验证评分", 10.3, 1.5, ORANGE)

    add_text(slide, "输出格式", 5.3, 3.4, 4.0, 0.35, size=20, color=INK, bold=True)
    add_text(slide, '{"cat": 2, "duck": 1, "deer": 1}', 5.3, 3.85, 6.0, 0.45,
             size=22, color=BLUE, bold=True)
    add_bullets(slide, [
        "单张图片输入 → 干净的字典输出",
        "仅保留支持的动物类别",
        "经检测过滤和重复抑制后计数"
    ], 5.3, 4.6, 6.5, size=14)
    add_footer(slide, 1)
    return slide


# ============================================================
# SLIDE 2: 任务要求
# ============================================================
def slide02(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, PAPER)
    add_kicker(slide, "任务要求")
    add_title_block(slide, "结构化的检测到字典转换问题",
                    "检测器需同时满足视觉识别和严格的JSON格式要求")

    steps = [
        (1, "检测", "发现并定位每个动物实例", BLUE),
        (2, "过滤", "仅保留20个允许的动物类别", CYAN),
        (3, "抑制", "移除重复预测框后再计数", ORANGE),
        (4, "计数", "按类别标签分组统计", GREEN),
        (5, "导出", "将批量结果写入JSON字典", RED),
    ]
    for num, title, desc, color in steps:
        add_step(slide, num, title, desc, 0.6 + (num - 1) * 2.45, 2.6, color)

    add_text(slide, "评分重点", 0.6, 4.5, 3.5, 0.35, size=20, color=INK, bold=True)
    add_bullets(slide, [
        "类别召回率：是否每个真实物种都被检测到？",
        "计数准确性：每个类别的计数是否正确？",
        "误报惩罚：每个额外的错误类别扣5分"
    ], 0.6, 5.0, 11.0, size=14)

    add_text(slide, "评估指标", 0.6, 6.1, 3.5, 0.35, size=16, color=INK, bold=True)
    add_text(slide, "格式有效性 | 类别精确度 | 类别召回率 | 计数准确度 | 整体响应质量",
             0.6, 6.45, 12.0, 0.35, size=10, color=MUTED)
    add_footer(slide, 2)
    return slide


# ============================================================
# SLIDE 3: 数据集构建
# ============================================================
def slide03(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, PAPER)
    add_kicker(slide, "数据集构建")
    add_title_block(slide, "分阶段扩展训练数据，保留完整20类词汇表",
                    "每个微调阶段均保留20类输出头")

    stages = [
        ("1. 合并", "xfz + yzy 数据集", "将标签转换为JSON格式", BLUE),
        ("2. 训练", "animal_dataset1 (20类)", "20类YOLO基线模型", CYAN),
        ("3. 优化", "animal_dataset_5sp_v1", "针对狼/鹅/马/牛/猴弱类别", ORANGE),
        ("4. 继续", "animal_dataset_6sp_v1", "从最佳模型继续优化", GREEN),
    ]
    for i, (label, data, desc, color) in enumerate(stages):
        top = 2.5 + i * 0.9
        add_rounded(slide, 0.6, top, 11.6, 0.65, WHITE, LINE)
        add_rect(slide, 0.6, top, 0.06, 0.65, color)
        add_text(slide, label, 0.85, top + 0.15, 2.8, 0.3, size=14, color=color, bold=True)
        add_text(slide, data, 3.8, top + 0.15, 3.5, 0.3, size=14, color=INK, bold=True)
        add_text(slide, desc, 7.5, top + 0.15, 4.5, 0.3, size=12, color=MUTED)

    add_text(slide, "标签格式", 0.6, 6.3, 2.5, 0.3, size=16, color=INK, bold=True)
    add_text(slide, "LabelMe风格JSON（标注审查）→ YOLO TXT格式（训练）",
             3.3, 6.3, 9.0, 0.5, size=12, color=MUTED)
    add_footer(slide, 3)
    return slide


# ============================================================
# SLIDE 4: 模型流程
# ============================================================
def slide04(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, PAPER)
    add_kicker(slide, "模型流程")
    add_title_block(slide, "YOLO推理与任务特定后处理相结合",
                    "后处理层至关重要，因为评分目标是字典而非边界框列表")

    inf_steps = [
        (1, "YOLO11s", "COCO预训练检测器", BLUE),
        (2, "微调", "dataset1 → 5sp → 6sp 三阶段", CYAN),
        (3, "白名单", "仅保留20个支持标签", GREEN),
        (4, "NMS+", "平衡式近距框抑制", ORANGE),
        (5, "JSON", "仅输出正计数", RED),
    ]
    for num, title, desc, color in inf_steps:
        add_step(slide, num, title, desc, 0.6 + (num - 1) * 2.45, 2.5, color)

    add_text(slide, "训练配置", 0.6, 4.3, 3.5, 0.35, size=18, color=INK, bold=True)
    add_bullets(slide, [
        "CUDA环境：NVIDIA RTX 3070 Ti 笔记本GPU",
        "图像尺寸640，批次大小8，冻结前10层，早停耐心值15",
        "数据增强：HSV抖动、平移、缩放、翻转、马赛克、轻量mixup",
        "优化器：AdamW，初始学习率0.0005，权重衰减0.0005，预热3轮"
    ], 0.6, 4.75, 11.5, size=13)
    add_footer(slide, 4)
    return slide


# ============================================================
# SLIDE 5: 训练结果
# ============================================================
def slide05(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, PAPER)
    add_kicker(slide, "训练结果")
    add_title_block(slide, "最终微调快速收敛，第2轮被选为最佳模型",
                    "验证指标保持高位，后期轮次无持续提升")

    curves = PROJECT_ROOT / "runs" / "detect" / "runs" / "detect" / "animal_dataset_6sp_20cls_finetune" / "results.png"
    if curves.exists():
        slide.shapes.add_picture(str(curves), Inches(0.5), Inches(2.3), Inches(7.5), Inches(4.0))

    add_metric(slide, "0.963", "精确率", 8.6, 2.5, BLUE)
    add_metric(slide, "0.977", "召回率", 10.9, 2.5, GREEN)
    add_metric(slide, "0.978", "mAP50", 8.6, 4.3, ORANGE)
    add_metric(slide, "0.940", "mAP50-95", 10.9, 4.3, RED)

    add_text(slide, "最佳模型：runs/detect/.../animal_dataset_6sp_20cls_finetune/weights/best.pt",
             8.6, 6.2, 4.0, 0.6, size=11, color=MUTED)

    add_text(slide, "训练阶段", 0.5, 6.5, 3.0, 0.3, size=15, color=INK, bold=True)
    for i, (stage, data, base) in enumerate([
        ("阶段1", "animal_dataset1", "COCO预训练权重"),
        ("阶段2", "animal_dataset_5sp_v1", "阶段1 best.pt"),
        ("阶段3", "animal_dataset_6sp_v1", "阶段2 best.pt"),
    ]):
        left = 0.5 + i * 4.1
        add_text(slide, f"{stage}: {data}", left, 6.8, 3.8, 0.2, size=10, color=BLUE, bold=True)
        add_text(slide, f"基于 {base}", left, 7.0, 3.8, 0.2, size=9, color=MUTED)
    add_footer(slide, 5)
    return slide


# ============================================================
# SLIDE 6: 验证结果
# ============================================================
def slide06(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, PAPER)
    add_kicker(slide, "验证结果")
    add_title_block(slide, "平衡式抑制策略：平均验证评分 90.94/100",
                    "剩余错误集中在易混淆的禽类和犬科类别")

    add_metric(slide, "90.94", "平均分 /100", 0.6, 2.5, ORANGE)
    add_metric(slide, "9/15", "满分图片", 3.1, 2.5, GREEN)
    add_metric(slide, "15", "验证集图片", 5.6, 2.5, BLUE)

    add_text(slide, "逐图评分", 0.6, 4.3, 3.0, 0.35, size=16, color=INK, bold=True)
    scores = [
        ("eval_001384", 100), ("eval_001388", 100), ("eval_001431", 100),
        ("val_4", 100), ("val_5", 100), ("val_6", 100),
        ("val_7", 100), ("val_8", 100), ("val_10", 100),
        ("eval_001391", 83.33), ("val_1", 75),
        ("val_2", 61.67), ("val_9", 61.67),
        ("eval_001402", 50), ("val_3", 50),
    ]
    for i, (name, score) in enumerate(scores):
        col, row = i % 2, i // 2
        left, top = 0.6 + col * 6.2, 4.8 + row * 0.32
        bar_c = GREEN if score >= 100 else (ORANGE if score >= 70 else RED)
        add_bar(slide, name, score, 100, left, top, bar_c)

    add_text(slide, "主要错误类型", 8.4, 2.6, 4.0, 0.3, size=16, color=INK, bold=True)
    add_bullets(slide, [
        "漏检（如 val_1 中的狐狸）",
        "类别混淆：鸭 ↔ 鹅",
        "类别混淆：狐狸/狼/狗",
        "额外的误报类别",
        "遮挡场景中的计数错误"
    ], 8.4, 3.1, 4.0, size=12)
    add_footer(slide, 6)
    return slide


# ============================================================
# SLIDE 7: 可视化输出
# ============================================================
def slide07(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, PAPER)
    add_kicker(slide, "可视化输出")
    add_title_block(slide, "验证集图片标注了边界框和置信度分数",
                    "显示框和JSON计数使用相同的后处理逻辑")

    anno = PROJECT_ROOT / "outputs" / "val_set_annotated_6sp_20cls_finetuned_suppressed_balanced" / "annotated" / "eval_001384.png"
    if anno.exists():
        slide.shapes.add_picture(str(anno), Inches(0.6), Inches(2.3), Inches(5.8), Inches(4.3))

    add_text(slide, "后处理规则", 7.0, 2.5, 5.0, 0.35, size=18, color=INK, bold=True)
    add_bullets(slide, [
        "同类别：IoU高或中心距近且尺寸相似时合并",
        "跨类别：仅在很可能覆盖同一目标时合并（更严格阈值）",
        "在重复检测中保留置信度最高的",
        "平衡模式：减少重复但不激进删除邻近动物"
    ], 7.0, 3.0, 5.5, size=12)

    add_text(slide, "输出目录", 7.0, 5.5, 3.0, 0.3, size=14, color=BLUE, bold=True)
    add_text(slide, "outputs/val_set_annotated_6sp_20cls_finetuned_suppressed_balanced/\n"
             "  annotated/  ← 带框图片\n"
             "  detections_with_boxes.json  ← 完整检测元数据",
             7.0, 5.85, 5.5, 0.8, size=10, color=MUTED)
    add_footer(slide, 7)
    return slide


# ============================================================
# SLIDE 8: 分析与挑战
# ============================================================
def slide08(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, PAPER)
    add_kicker(slide, "分析与挑战")
    add_title_block(slide, "深入分析模型行为和遗留问题")

    for i, (title, color, items) in enumerate([
        ("优势", GREEN, [
            "大型特征明显动物检测良好（斑马、长颈鹿等）",
            "高 mAP50 (0.978) 和 mAP50-95 (0.940)",
            "三阶段微调改善了弱类别表现",
            "干净的JSON输出格式，无自然语言解释"
        ]),
        ("局限", RED, [
            "鸭与鹅混淆（体型和颜色相似）",
            "狐狸/狼/狗混淆（纹理和姿态相似）",
            "遮挡和密集排列动物处理困难",
            "合成场景到真实场景可能存在域差距"
        ]),
        ("解决方案", BLUE, [
            "自定义后处理进行重复框抑制",
            "针对弱类别数据集定向微调",
            "平衡阈值避免激进过滤",
            "全程保持20类输出头"
        ]),
    ]):
        left = 0.5 + i * 4.2
        add_rounded(slide, left, 2.2, 3.9, 2.5, WHITE, LINE)
        add_rect(slide, left, 2.2, 3.9, 0.06, color)
        add_text(slide, title, left + 0.2, 2.4, 3.5, 0.3, size=18, color=color, bold=True)
        add_bullets(slide, items, left + 0.2, 2.85, 3.5, size=11)

    # Key findings
    add_rounded(slide, 0.5, 5.0, 12.3, 1.8, RGBColor(0xD9, 0xEA, 0xF7), BLUE)
    add_text(slide, "关键发现", 0.7, 5.15, 4.0, 0.3, size=16, color=BLUE, bold=True)
    add_bullets(slide, [
        "微调时保持20类输出避免了词汇丢失，同时改善了弱类别表现",
        "平衡后处理（同类别IoU=0.40，跨类别IoU=0.62）减少重复但未损害召回率",
        "评分规则对误报类别惩罚严厉（每类5分）→ 需要保守的阈值调优"
    ], 0.7, 5.55, 11.8, size=12)
    add_footer(slide, 8)
    return slide


# ============================================================
# SLIDE 9: 团队贡献
# ============================================================
def slide09(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, PAPER)
    add_kicker(slide, "团队贡献")
    add_title_block(slide, "每位成员在项目中的角色与贡献")

    members = [
        ("成员1", "待填写", [
            "数据集收集与合并",
            "LabelMe JSON标注流程",
            "验证数据组织",
            "数据预处理脚本"
        ], BLUE),
        ("成员2", "待填写", [
            "YOLO模型训练与微调",
            "CUDA环境配置与优化",
            "模型指标跟踪与选型",
            "超参数调优"
        ], GREEN),
        ("成员3", "待填写", [
            "评估与推理脚本",
            "重复框后处理算法",
            "报告、PPT与文档撰写",
            "错误分析与验证评分"
        ], ORANGE),
    ]
    for i, (name, sid, tasks, color) in enumerate(members):
        left = 0.5 + i * 4.2
        add_rounded(slide, left, 2.3, 3.9, 4.3, WHITE, LINE)
        add_rect(slide, left, 2.3, 3.9, 0.06, color)
        add_text(slide, name, left + 0.2, 2.5, 3.5, 0.35, size=20, color=color, bold=True)
        add_text(slide, f"学号: {sid}", left + 0.2, 2.85, 3.5, 0.25, size=12, color=MUTED)
        add_rect(slide, left + 0.2, 3.2, 3.5, 0.01, LINE)
        add_bullets(slide, tasks, left + 0.2, 3.35, 3.5, size=13)
    add_footer(slide, 9)
    return slide


# ============================================================
# SLIDE 10: 结论
# ============================================================
def slide10(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, DARK_BG)

    add_text(slide, "结论", 0.6, 0.8, 6.0, 0.7, size=38, color=WHITE, bold=True)
    add_text(slide, "可复现的基于YOLO的动物计数系统，涵盖数据准备、微调、验证评分、JSON导出和可视化标注",
             0.7, 1.7, 10.0, 0.65, size=18, color=RGBColor(0xDC, 0xE8, 0xF8))

    cards = [
        ("成果", GREEN, [
            "20个动物类别 mAP50 达 0.978",
            "15张验证集评分 90.94/100",
            "干净的字典格式JSON输出",
            "完整的代码+报告+PPT交付物"
        ]),
        ("局限", ORANGE, [
            "鸭/鹅和狐狸/狼/狗的混淆仍然存在",
            "遮挡和密集排列导致计数错误",
            "合成到真实场景可能存在域差距"
        ]),
        ("后续工作", BLUE, [
            "增加弱类别的人工验证难例",
            "以更严格审查继续定向微调",
            "探索测试时增强提升鲁棒性",
            "改善鸭/鹅区分特征"
        ]),
    ]
    for i, (title, color, items) in enumerate(cards):
        left = 0.5 + i * 4.2
        add_rounded(slide, left, 2.8, 3.9, 3.8, RGBColor(0x14, 0x2A, 0x46), RGBColor(0x35, 0x5B, 0x88))
        add_text(slide, title, left + 0.2, 3.0, 3.5, 0.35, size=20, color=color, bold=True)
        txBox = slide.shapes.add_textbox(Inches(left + 0.2), Inches(3.5), Inches(3.5), Inches(2.8))
        tf = txBox.text_frame; tf.word_wrap = True
        for j, item in enumerate(items):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.text = f"▸ {item}"
            p.font.size = Pt(13); p.font.color.rgb = WHITE; p.font.name = ZH_FONT
            p.space_after = Pt(6)

    add_text(slide, "最终交付物：报告 | PPT | 源代码 | 评估JSON | 标注图片",
             0.7, 6.9, 11.5, 0.35, size=13, color=RGBColor(0xBB, 0xD2, 0xF0))
    add_footer(slide, 10)
    return slide


def main():
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide01(prs); slide02(prs); slide03(prs); slide04(prs); slide05(prs)
    slide06(prs); slide07(prs); slide08(prs); slide09(prs); slide10(prs)

    prs.save(str(PPT_PATH))
    print(f"Chinese PPT saved: {PPT_PATH} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
