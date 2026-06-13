"""
Generate English-only PPT presentation for the Animal Detection project.
"""
from __future__ import annotations

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

PROJECT_ROOT = Path(__file__).resolve().parents[1]
OUTPUT_DIR = PROJECT_ROOT / "outputs" / "final_deliverables"
PPT_PATH = OUTPUT_DIR / "animal_detection_project_presentation_en.pptx"

# Colors
DARK_BG = RGBColor(0x0E, 0x1A, 0x2B)
BLUE = RGBColor(0x24, 0x5A, 0x9A)
CYAN = RGBColor(0x32, 0xA6, 0xB8)
GREEN = RGBColor(0x4B, 0x9B, 0x68)
ORANGE = RGBColor(0xE2, 0x8A, 0x32)
RED = RGBColor(0xC9, 0x4F, 0x4F)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x17, 0x20, 0x33)
MUTED = RGBColor(0x5E, 0x6A, 0x7D)
PAPER = RGBColor(0xF7, 0xF9, 0xFC)
LINE = RGBColor(0xD8, 0xE2, 0xEF)


def add_bg(slide, color=PAPER):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    return bg


def add_text(slide, text, left, top, width, height, size=18, color=INK,
             bold=False, align=PP_ALIGN.LEFT, font="Arial"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.alignment = align
    return txBox


def add_multiline(slide, lines, left, top, width, size=18, color=INK, spacing=1.2):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(6))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Arial"
        p.space_after = Pt(4)
    return txBox


def add_bullets(slide, items, left, top, width, size=15, color=INK):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(5.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"▸ {item}"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Arial"
        p.space_after = Pt(6)
    return txBox


def add_rect(slide, left, top, width, height, fill, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_rounded(slide, left, top, width, height, fill, line=LINE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(1)
    return shape


def add_footer(slide, n, total=10):
    add_rect(slide, 0.5, 7.05, 12.3, 0.01, LINE)
    add_text(slide, "YOLO11s Animal Detection & Counting | Deep Learning Final Project",
             0.5, 7.1, 9.5, 0.35, size=9, color=MUTED)
    add_text(slide, f"{n}/{total}", 11.5, 7.1, 0.8, 0.35, size=11, color=BLUE,
             bold=True, align=PP_ALIGN.RIGHT)


def add_kicker(slide, text):
    marker = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    Inches(0.5), Inches(0.45), Inches(0.15), Inches(0.15))
    marker.fill.solid()
    marker.fill.fore_color.rgb = CYAN
    marker.line.fill.background()
    add_text(slide, text, 0.8, 0.35, 6.0, 0.35, size=11, color=BLUE, bold=True)


def add_title_block(slide, title, subtitle=""):
    add_text(slide, title, 0.5, 0.85, 12.0, 0.8, size=30, color=INK, bold=True)
    if subtitle:
        add_text(slide, subtitle, 0.55, 1.55, 11.0, 0.4, size=15, color=MUTED)


def add_metric(slide, value, label, left, top, color=BLUE):
    add_rounded(slide, left, top, 2.2, 1.5, WHITE, LINE)
    add_text(slide, value, left + 0.2, top + 0.15, 1.8, 0.45, size=26, color=color, bold=True)
    add_text(slide, label, left + 0.2, top + 0.65, 1.8, 0.5, size=11, color=MUTED)


def add_step(slide, num, title, desc, left, top, color=BLUE):
    add_rounded(slide, left, top, 2.0, 1.35, WHITE, LINE)
    add_rect(slide, left, top, 2.0, 0.08, color)
    add_text(slide, f"Step {num}", left + 0.15, top + 0.18, 1.7, 0.25, size=10, color=color, bold=True)
    add_text(slide, title, left + 0.15, top + 0.45, 1.7, 0.28, size=15, color=INK, bold=True)
    add_text(slide, desc, left + 0.15, top + 0.78, 1.7, 0.45, size=9, color=MUTED)


def add_bar(slide, label, value, max_val, left, top, color=GREEN):
    bw = 5.0
    add_text(slide, label, left, top, 2.2, 0.28, size=12, color=INK)
    add_rect(slide, left + 2.4, top + 0.06, bw, 0.16, RGBColor(0xE8, 0xEE, 0xF6))
    if value > 0:
        add_rect(slide, left + 2.4, top + 0.06, bw * value / max_val, 0.16, color)
    add_text(slide, f"{value:.1f}" if isinstance(value, float) else str(value),
             left + 2.4 + bw + 0.15, top, 0.6, 0.28, size=12, color=color, bold=True)


# ============================================================
# SLIDE 1: COVER
# ============================================================
def slide01(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, PAPER)
    add_rect(slide, 0, 0, 4.3, 7.5, DARK_BG)

    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(3.3), Inches(1.8))
    tf = txBox.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = "Animal Detection\nand Counting"
    p.font.size = Pt(36); p.font.color.rgb = WHITE; p.font.bold = True; p.font.name = "Arial"
    p2 = tf.add_paragraph(); p2.text = "YOLO11s + Custom 20-Class Fine-Tuning"
    p2.font.size = Pt(16); p2.font.color.rgb = RGBColor(0xBB, 0xD2, 0xF0); p2.font.name = "Arial"

    add_metric(slide, "20", "Animal Categories", 5.3, 1.5, BLUE)
    add_metric(slide, "0.978", "Final mAP50", 7.8, 1.5, GREEN)
    add_metric(slide, "90.94", "Validation Score", 10.3, 1.5, ORANGE)

    add_text(slide, "Output Contract", 5.3, 3.4, 4.0, 0.35, size=20, color=INK, bold=True)
    add_text(slide, '{"cat": 2, "duck": 1, "deer": 1}', 5.3, 3.85, 6.0, 0.45,
             size=22, color=BLUE, bold=True)
    add_bullets(slide, [
        "One image in → one clean dictionary out",
        "Only supported animal categories are retained",
        "Counts generated after filtering & duplicate suppression"
    ], 5.3, 4.6, 6.5, size=14)
    add_footer(slide, 1)
    return slide


# ============================================================
# SLIDE 2: TASK REQUIREMENTS
# ============================================================
def slide02(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, PAPER)
    add_kicker(slide, "TASK REQUIREMENTS")
    add_title_block(slide, "Structured Detection-to-Dictionary Problem",
                    "The detector must satisfy both visual recognition and strict JSON formatting.")

    steps = [
        (1, "Detect", "Find & localize every animal instance", BLUE),
        (2, "Filter", "Keep only 20 allowed animal categories", CYAN),
        (3, "Suppress", "Remove duplicate predictions before counting", ORANGE),
        (4, "Count", "Group remaining boxes by class label", GREEN),
        (5, "Export", "Write batch results as JSON dictionaries", RED),
    ]
    for num, title, desc, color in steps:
        add_step(slide, num, title, desc, 0.6 + (num - 1) * 2.45, 2.6, color)

    add_text(slide, "Scoring Focus", 0.6, 4.5, 3.5, 0.35, size=20, color=INK, bold=True)
    add_bullets(slide, [
        "Category recall: does every true species appear?",
        "Counting accuracy: is each class count correct?",
        "False positives: each extra wrong category costs 5 points"
    ], 0.6, 5.0, 11.0, size=14)

    add_text(slide, "Evaluation Metrics", 0.6, 6.1, 3.5, 0.35, size=16, color=INK, bold=True)
    add_text(slide, "Format Validity | Category Precision | Category Recall | Counting Accuracy | Overall Response Quality",
             0.6, 6.45, 12.0, 0.35, size=10, color=MUTED)
    add_footer(slide, 2)
    return slide


# ============================================================
# SLIDE 3: DATASET
# ============================================================
def slide03(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, PAPER)
    add_kicker(slide, "DATASET CONSTRUCTION")
    add_title_block(slide, "Training Data Expanded in Stages",
                    "Every fine-tuning stage preserved the 20-class output head.")

    stages = [
        ("1. Merge", "xfz + yzy datasets", "Convert labels to JSON format", BLUE),
        ("2. Train", "animal_dataset1 (20cls)", "20-class YOLO baseline model", CYAN),
        ("3. Improve", "animal_dataset_5sp_v1", "Target: wolf/goose/horse/cow/monkey", ORANGE),
        ("4. Continue", "animal_dataset_6sp_v1", "Optimization from best model", GREEN),
    ]
    for i, (label, data, desc, color) in enumerate(stages):
        top = 2.5 + i * 0.9
        add_rounded(slide, 0.6, top, 11.6, 0.65, WHITE, LINE)
        add_rect(slide, 0.6, top, 0.06, 0.65, color)
        add_text(slide, label, 0.85, top + 0.15, 2.8, 0.3, size=14, color=color, bold=True)
        add_text(slide, data, 3.8, top + 0.15, 3.5, 0.3, size=14, color=INK, bold=True)
        add_text(slide, desc, 7.5, top + 0.15, 4.5, 0.3, size=12, color=MUTED)

    add_text(slide, "Label Format", 0.6, 6.3, 2.5, 0.3, size=16, color=INK, bold=True)
    add_text(slide, "LabelMe-style JSON for annotation review → YOLO TXT format for training",
             3.3, 6.3, 9.0, 0.5, size=12, color=MUTED)
    add_footer(slide, 3)
    return slide


# ============================================================
# SLIDE 4: MODEL PIPELINE
# ============================================================
def slide04(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, PAPER)
    add_kicker(slide, "MODEL PIPELINE")
    add_title_block(slide, "YOLO Inference + Task-Specific Post-Processing",
                    "Post-processing is critical because grading targets a dictionary, not a box list.")

    inf_steps = [
        (1, "YOLO11s", "COCO pre-trained detector", BLUE),
        (2, "Fine-tune", "dataset1 → 5sp → 6sp stages", CYAN),
        (3, "Whitelist", "20 supported labels only", GREEN),
        (4, "NMS+", "Balanced close-box suppression", ORANGE),
        (5, "JSON", "Positive counts only output", RED),
    ]
    for num, title, desc, color in inf_steps:
        add_step(slide, num, title, desc, 0.6 + (num - 1) * 2.45, 2.5, color)

    add_text(slide, "Training Configuration", 0.6, 4.3, 3.5, 0.35, size=18, color=INK, bold=True)
    add_bullets(slide, [
        "CUDA: NVIDIA RTX 3070 Ti Laptop GPU",
        "imgsz=640, batch=8, freeze=10, patience=15",
        "Augmentations: HSV jitter, translation, scale, flip, mosaic, light mixup",
        "Optimizer: AdamW, lr0=0.0005, weight_decay=0.0005, warmup_epochs=3"
    ], 0.6, 4.75, 11.5, size=13)
    add_footer(slide, 4)
    return slide


# ============================================================
# SLIDE 5: TRAINING RESULTS
# ============================================================
def slide05(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, PAPER)
    add_kicker(slide, "TRAINING RESULTS")
    add_title_block(slide, "Final Fine-Tuning Converged Quickly; Epoch 2 Selected as Best",
                    "Validation metrics remained high while later epochs showed no sustained improvement.")

    curves = PROJECT_ROOT / "runs" / "detect" / "runs" / "detect" / "animal_dataset_6sp_20cls_finetune" / "results.png"
    if curves.exists():
        slide.shapes.add_picture(str(curves), Inches(0.5), Inches(2.3), Inches(7.5), Inches(4.0))

    add_metric(slide, "0.963", "Precision", 8.6, 2.5, BLUE)
    add_metric(slide, "0.977", "Recall", 10.9, 2.5, GREEN)
    add_metric(slide, "0.978", "mAP50", 8.6, 4.3, ORANGE)
    add_metric(slide, "0.940", "mAP50-95", 10.9, 4.3, RED)

    add_text(slide, "Best: runs/detect/.../animal_dataset_6sp_20cls_finetune/weights/best.pt",
             8.6, 6.2, 4.0, 0.6, size=11, color=MUTED)

    add_text(slide, "Training Stages", 0.5, 6.5, 3.0, 0.3, size=15, color=INK, bold=True)
    for i, (stage, data, base) in enumerate([
        ("Stage 1", "animal_dataset1", "COCO pre-trained"),
        ("Stage 2", "animal_dataset_5sp_v1", "Stage 1 best.pt"),
        ("Stage 3", "animal_dataset_6sp_v1", "Stage 2 best.pt"),
    ]):
        left = 0.5 + i * 4.1
        add_text(slide, f"{stage}: {data}", left, 6.8, 3.8, 0.2, size=10, color=BLUE, bold=True)
        add_text(slide, f"From {base}", left, 7.0, 3.8, 0.2, size=9, color=MUTED)
    add_footer(slide, 5)
    return slide


# ============================================================
# SLIDE 6: VALIDATION
# ============================================================
def slide06(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, PAPER)
    add_kicker(slide, "VALIDATION RESULTS")
    add_title_block(slide, "Balanced Suppression: 90.94/100 Average Score",
                    "Remaining errors concentrated in confusing bird & canine-like classes.")

    add_metric(slide, "90.94", "Avg Score /100", 0.6, 2.5, ORANGE)
    add_metric(slide, "9/15", "Perfect Images", 3.1, 2.5, GREEN)
    add_metric(slide, "15", "Validation Set", 5.6, 2.5, BLUE)

    add_text(slide, "Per-Image Scores", 0.6, 4.3, 3.0, 0.35, size=16, color=INK, bold=True)
    scores = [
        ("eval_001384", 100), ("eval_001388", 100), ("eval_001431", 100),
        ("val_4", 100), ("val_5", 100), ("val_6", 100),
        ("val_7", 100), ("val_8", 100), ("val_10", 100),
        ("eval_001391", 83.33), ("val_1", 75),
        ("val_2", 61.67), ("val_9", 61.67),
        ("eval_001402", 50), ("val_3", 50),
    ]
    for i, (name, score) in enumerate(scores):
        col, row = i % 2, i // 2
        left, top = 0.6 + col * 6.2, 4.8 + row * 0.32
        bar_c = GREEN if score >= 100 else (ORANGE if score >= 70 else RED)
        add_bar(slide, name, score, 100, left, top, bar_c)

    add_text(slide, "Main Error Types", 8.4, 2.6, 4.0, 0.3, size=16, color=INK, bold=True)
    add_bullets(slide, [
        "Missed detections (e.g., fox in val_1)",
        "Class confusion: duck ↔ goose",
        "Class confusion: fox/wolf/dog",
        "Extra false positive categories",
        "Count errors in occluded scenes"
    ], 8.4, 3.1, 4.0, size=12)
    add_footer(slide, 6)
    return slide


# ============================================================
# SLIDE 7: VISUAL OUTPUT
# ============================================================
def slide07(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, PAPER)
    add_kicker(slide, "VISUAL OUTPUT")
    add_title_block(slide, "Annotated Validation Images with Boxes & Confidence Scores",
                    "Same post-processing logic used for displayed boxes and JSON counts.")

    anno = PROJECT_ROOT / "outputs" / "val_set_annotated_6sp_20cls_finetuned_suppressed_balanced" / "annotated" / "eval_001384.png"
    if anno.exists():
        slide.shapes.add_picture(str(anno), Inches(0.6), Inches(2.3), Inches(5.8), Inches(4.3))

    add_text(slide, "Post-Processing Rules", 7.0, 2.5, 5.0, 0.35, size=18, color=INK, bold=True)
    add_bullets(slide, [
        "Same-class: merge when IoU high or centers close with similar size",
        "Cross-class: merge only when likely covering same object (stricter)",
        "Keep highest-confidence detection among duplicates",
        "Balanced mode: reduce duplicates without aggressive deletion"
    ], 7.0, 3.0, 5.5, size=12)

    add_text(slide, "Output Directory", 7.0, 5.5, 3.0, 0.3, size=14, color=BLUE, bold=True)
    add_text(slide, "outputs/val_set_annotated_6sp_20cls_finetuned_suppressed_balanced/\n"
             "  annotated/  ← images with boxes\n"
             "  detections_with_boxes.json  ← full metadata",
             7.0, 5.85, 5.5, 0.8, size=10, color=MUTED)
    add_footer(slide, 7)
    return slide


# ============================================================
# SLIDE 8: ANALYSIS
# ============================================================
def slide08(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, PAPER)
    add_kicker(slide, "ANALYSIS & CHALLENGES")
    add_title_block(slide, "Deep Dive into Model Behavior and Remaining Issues")

    for i, (title, color, items) in enumerate([
        ("Strengths", GREEN, [
            "Large distinctive animals detected excellently",
            "High mAP50 (0.978) and mAP50-95 (0.940)",
            "Three-stage fine-tuning improved weak classes",
            "Clean JSON output format, no NL explanations"
        ]),
        ("Limitations", RED, [
            "Duck ↔ goose confusion (similar body shape & color)",
            "Fox/wolf/dog confusion (similar texture & pose)",
            "Occlusion & tightly grouped animals",
            "Synthetic scene generalization gap"
        ]),
        ("Solutions", BLUE, [
            "Custom post-processing for duplicate suppression",
            "Targeted fine-tuning on weak-class datasets",
            "Balanced thresholds to avoid aggressive filtering",
            "Maintained 20-class head throughout all stages"
        ]),
    ]):
        left = 0.5 + i * 4.2
        add_rounded(slide, left, 2.2, 3.9, 2.5, WHITE, LINE)
        add_rect(slide, left, 2.2, 3.9, 0.06, color)
        add_text(slide, title, left + 0.2, 2.4, 3.5, 0.3, size=18, color=color, bold=True)
        add_bullets(slide, items, left + 0.2, 2.85, 3.5, size=11)

    # Key findings
    add_rounded(slide, 0.5, 5.0, 12.3, 1.8, RGBColor(0xD9, 0xEA, 0xF7), BLUE)
    add_text(slide, "Key Findings", 0.7, 5.15, 4.0, 0.3, size=16, color=BLUE, bold=True)
    add_bullets(slide, [
        "Maintaining 20-class output during fine-tuning prevented vocabulary loss while improving weak categories",
        "Balanced post-processing (IoU=0.40 same-class, 0.62 cross-class) reduced duplicates without hurting recall",
        "Scoring rule penalizes false positive categories (5 pts each) → conservative threshold tuning needed"
    ], 0.7, 5.55, 11.8, size=12)
    add_footer(slide, 8)
    return slide


# ============================================================
# SLIDE 9: CONTRIBUTIONS
# ============================================================
def slide09(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, PAPER)
    add_kicker(slide, "TEAM CONTRIBUTIONS")
    add_title_block(slide, "Each Member's Role and Contribution to the Project")

    members = [
        ("Member 1", "TBD / 待填写", [
            "Dataset collection & merging",
            "LabelMe JSON annotation workflow",
            "Validation data organization",
            "Data preprocessing scripts"
        ], BLUE),
        ("Member 2", "TBD / 待填写", [
            "YOLO model training & fine-tuning",
            "CUDA environment setup & optimization",
            "Model metric tracking & selection",
            "Hyperparameter tuning"
        ], GREEN),
        ("Member 3", "TBD / 待填写", [
            "Evaluation & inference scripts",
            "Duplicate-box post-processing algorithm",
            "Report, PPT & documentation",
            "Error analysis & validation scoring"
        ], ORANGE),
    ]
    for i, (name, sid, tasks, color) in enumerate(members):
        left = 0.5 + i * 4.2
        add_rounded(slide, left, 2.3, 3.9, 4.3, WHITE, LINE)
        add_rect(slide, left, 2.3, 3.9, 0.06, color)
        add_text(slide, name, left + 0.2, 2.5, 3.5, 0.35, size=18, color=color, bold=True)
        add_text(slide, f"ID: {sid}", left + 0.2, 2.85, 3.5, 0.25, size=12, color=MUTED)
        add_rect(slide, left + 0.2, 3.2, 3.5, 0.01, LINE)
        add_bullets(slide, tasks, left + 0.2, 3.35, 3.5, size=12)
    add_footer(slide, 9)
    return slide


# ============================================================
# SLIDE 10: CONCLUSION
# ============================================================
def slide10(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, DARK_BG)

    add_text(slide, "Conclusion", 0.6, 0.8, 6.0, 0.7, size=38, color=WHITE, bold=True)
    add_text(slide, "A reproducible YOLO-based animal counting system covering data preparation, "
             "fine-tuning, validation scoring, JSON export, and visual annotation.",
             0.7, 1.7, 10.0, 0.65, size=18, color=RGBColor(0xDC, 0xE8, 0xF8))

    cards = [
        ("Achievements", GREEN, [
            "High mAP50 of 0.978 across 20 animal classes",
            "Validation score of 90.94/100 on 15-image set",
            "Clean dictionary-format JSON output",
            "Complete code + report + PPT deliverables"
        ]),
        ("Limitations", ORANGE, [
            "Duck/goose and fox/wolf/dog confusions persist",
            "Occlusion & tight grouping cause count errors",
            "Synthetic-to-real domain gap possible"
        ]),
        ("Next Steps", BLUE, [
            "Add manually verified hard examples for weak classes",
            "Continue targeted fine-tuning with stricter review",
            "Explore test-time augmentation for robustness",
            "Improve duck/goose discrimination features"
        ]),
    ]
    for i, (title, color, items) in enumerate(cards):
        left = 0.5 + i * 4.2
        add_rounded(slide, left, 2.8, 3.9, 3.8, RGBColor(0x14, 0x2A, 0x46), RGBColor(0x35, 0x5B, 0x88))
        add_text(slide, title, left + 0.2, 3.0, 3.5, 0.35, size=18, color=color, bold=True)
        txBox = slide.shapes.add_textbox(Inches(left + 0.2), Inches(3.5), Inches(3.5), Inches(2.8))
        tf = txBox.text_frame; tf.word_wrap = True
        for j, item in enumerate(items):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.text = f"▸ {item}"
            p.font.size = Pt(13); p.font.color.rgb = WHITE; p.font.name = "Arial"
            p.space_after = Pt(6)

    add_text(slide, "Final Deliverables: Report | PPT | Source Code | Evaluation JSON | Annotated Images",
             0.7, 6.9, 11.5, 0.35, size=13, color=RGBColor(0xBB, 0xD2, 0xF0))
    add_footer(slide, 10)
    return slide


def main():
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide01(prs); slide02(prs); slide03(prs); slide04(prs); slide05(prs)
    slide06(prs); slide07(prs); slide08(prs); slide09(prs); slide10(prs)

    prs.save(str(PPT_PATH))
    print(f"English PPT saved: {PPT_PATH} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
