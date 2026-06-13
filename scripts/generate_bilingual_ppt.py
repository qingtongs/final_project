"""
Generate bilingual (Chinese/English) PPT presentation for the Animal Detection project.
"""
from __future__ import annotations

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

PROJECT_ROOT = Path(__file__).resolve().parents[1]
OUTPUT_DIR = PROJECT_ROOT / "outputs" / "final_deliverables"
PPT_PATH = OUTPUT_DIR / "animal_detection_project_presentation.pptx"

# Color palette
DARK_BG = RGBColor(0x0E, 0x1A, 0x2B)
BLUE = RGBColor(0x24, 0x5A, 0x9A)
CYAN = RGBColor(0x32, 0xA6, 0xB8)
GREEN = RGBColor(0x4B, 0x9B, 0x68)
ORANGE = RGBColor(0xE2, 0x8A, 0x32)
RED = RGBColor(0xC9, 0x4F, 0x4F)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x17, 0x20, 0x33)
MUTED = RGBColor(0x5E, 0x6A, 0x7D)
PAPER = RGBColor(0xF7, 0xF9, 0xFC)
LINE = RGBColor(0xD8, 0xE2, 0xEF)
HEADING_BLUE = RGBColor(0x1F, 0x4E, 0x79)
LIGHT_BLUE_BG = RGBColor(0xD9, 0xEA, 0xF7)


def add_bg(slide, color=PAPER):
    """Add solid background rectangle."""
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    return bg


def add_text_box(slide, text, left, top, width, height, font_size=18, color=INK,
                 bold=False, alignment=PP_ALIGN.LEFT, font_name="Arial"):
    """Add a text box with specified properties."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_multiline_text(slide, lines, left, top, width, height, font_size=18, color=INK,
                       bold_first=False, line_spacing=1.3):
    """Add a text box with multiple lines."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Arial"
        p.font.bold = (bold_first and i == 0)
        p.space_after = Pt(4)
    return txBox


def add_bilingual_text(slide, en_text, zh_text, left, top, width, height,
                       en_size=20, zh_size=16, en_color=INK, zh_color=MUTED):
    """Add bilingual (English + Chinese) text."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True

    p_en = tf.paragraphs[0]
    p_en.text = en_text
    p_en.font.size = Pt(en_size)
    p_en.font.color.rgb = en_color
    p_en.font.name = "Arial"
    p_en.font.bold = True
    p_en.space_after = Pt(4)

    p_zh = tf.add_paragraph()
    p_zh.text = zh_text
    p_zh.font.size = Pt(zh_size)
    p_zh.font.color.rgb = zh_color
    p_zh.font.name = "Microsoft YaHei"
    p_zh.space_after = Pt(2)

    return txBox


def add_shape_rect(slide, left, top, width, height, fill_color, line_color=None):
    """Add a rectangle shape."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_rounded_rect(slide, left, top, width, height, fill_color, line_color=LINE):
    """Add a rounded rectangle."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = line_color
    shape.line.width = Pt(1)
    return shape


def add_footer(slide, slide_num, total=10):
    """Add footer with slide number."""
    add_shape_rect(slide, 0.5, 7.05, 12.3, 0.01, LINE)
    add_text_box(slide,
                 "YOLO11s Animal Detection & Counting | Deep Learning Final Project | 深度学习期末项目",
                 0.5, 7.1, 9.0, 0.35, font_size=9, color=MUTED)
    add_text_box(slide, f"{slide_num}/{total}",
                 11.5, 7.1, 0.8, 0.35, font_size=11, color=BLUE, bold=True,
                 alignment=PP_ALIGN.RIGHT)


def add_kicker(slide, text_en, text_zh):
    """Add section kicker/tag."""
    marker = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    Inches(0.5), Inches(0.45), Inches(0.15), Inches(0.15))
    marker.fill.solid()
    marker.fill.fore_color.rgb = CYAN
    marker.line.fill.background()

    add_text_box(slide, f"{text_en}  |  {text_zh}",
                 0.8, 0.35, 6.0, 0.35, font_size=11, color=BLUE, bold=True)


def add_title(slide, en_title, zh_title, en_subtitle="", zh_subtitle=""):
    """Add bilingual title with optional subtitle."""
    add_bilingual_text(slide, en_title, zh_title, 0.5, 0.85, 12.0, 0.8,
                       en_size=30, zh_size=22, en_color=INK, zh_color=MUTED)
    if en_subtitle:
        add_text_box(slide, en_subtitle, 0.55, 1.55, 11.0, 0.4, font_size=15, color=MUTED)
    if zh_subtitle:
        add_text_box(slide, zh_subtitle, 0.55, 1.85, 11.0, 0.4, font_size=13, color=MUTED)


def add_metric_card(slide, value, label_en, label_zh, left, top, color=BLUE):
    """Add a metric card with value and bilingual label."""
    card = add_rounded_rect(slide, left, top, 2.2, 1.5, WHITE, LINE)
    add_text_box(slide, value, left + 0.2, top + 0.15, 1.8, 0.45,
                 font_size=26, color=color, bold=True)
    add_text_box(slide, label_en, left + 0.2, top + 0.65, 1.8, 0.3,
                 font_size=11, color=MUTED)
    add_text_box(slide, label_zh, left + 0.2, top + 0.9, 1.8, 0.3,
                 font_size=10, color=MUTED)


def add_pipeline_step(slide, step_num, en_label, zh_label, en_desc, zh_desc,
                      left, top, color=BLUE):
    """Add a pipeline step with bilingual content."""
    card = add_rounded_rect(slide, left, top, 2.0, 1.35, WHITE, LINE)
    # Top accent bar
    add_shape_rect(slide, left, top, 2.0, 0.08, color)
    # Step number
    add_text_box(slide, f"Step {step_num}", left + 0.15, top + 0.18, 1.7, 0.25,
                 font_size=10, color=color, bold=True)
    # English label
    add_text_box(slide, en_label, left + 0.15, top + 0.45, 1.7, 0.28,
                 font_size=15, color=INK, bold=True)
    # Chinese label
    add_text_box(slide, zh_label, left + 0.15, top + 0.72, 1.7, 0.22,
                 font_size=11, color=MUTED)
    # Description
    add_text_box(slide, en_desc, left + 0.15, top + 0.95, 1.7, 0.32,
                 font_size=9, color=MUTED)


def add_bullet_list(slide, items, left, top, width, font_size=15, color=INK):
    """Add bullet point list. Each item is (en_text, zh_text) tuple."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(5.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if isinstance(item, tuple):
            en_text, zh_text = item
            if i == 0:
                p_en = tf.paragraphs[0]
            else:
                p_en = tf.add_paragraph()
            p_en.text = f"▸ {en_text}"
            p_en.font.size = Pt(font_size)
            p_en.font.color.rgb = color
            p_en.font.name = "Arial"
            p_en.space_after = Pt(1)

            p_zh = tf.add_paragraph()
            p_zh.text = f"    {zh_text}"
            p_zh.font.size = Pt(font_size - 2)
            p_zh.font.color.rgb = MUTED
            p_zh.font.name = "Microsoft YaHei"
            p_zh.space_after = Pt(8)
        else:
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"▸ {item}"
            p.font.size = Pt(font_size)
            p.font.color.rgb = color
            p.font.name = "Arial"
            p.space_after = Pt(6)


def add_bar_chart_row(slide, label, value, max_val, left, top, color=GREEN):
    """Add a horizontal bar showing a score."""
    bar_width = 5.0
    add_text_box(slide, label, left, top, 2.2, 0.28, font_size=12, color=INK)
    # Background bar
    add_shape_rect(slide, left + 2.4, top + 0.06, bar_width, 0.16,
                   RGBColor(0xE8, 0xEE, 0xF6))
    # Value bar
    if value > 0:
        add_shape_rect(slide, left + 2.4, top + 0.06, bar_width * value / max_val, 0.16,
                       color)
    add_text_box(slide, f"{value:.1f}" if isinstance(value, float) else str(value),
                 left + 2.4 + bar_width + 0.15, top, 0.6, 0.28,
                 font_size=12, color=color, bold=True)


# ============================================================
# SLIDE 1: COVER / 封面
# ============================================================
def slide01_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    add_bg(slide, PAPER)

    # Left dark panel
    add_shape_rect(slide, 0, 0, 4.3, 7.5, DARK_BG)

    # Title on dark panel
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(3.3), Inches(1.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Animal Detection\nand Counting"
    p.font.size = Pt(36)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = "Arial"
    p2 = tf.add_paragraph()
    p2.text = "动物检测与计数系统"
    p2.font.size = Pt(24)
    p2.font.color.rgb = RGBColor(0xBB, 0xD2, 0xF0)
    p2.font.name = "Microsoft YaHei"

    # Subtitle
    add_text_box(slide, "YOLO11s + Custom 20-Class Fine-Tuning\n基于YOLO11s的20类自定义微调模型",
                 0.55, 3.2, 3.3, 0.7, font_size=14, color=RGBColor(0xBB, 0xD2, 0xF0))

    # Metrics on right side
    add_metric_card(slide, "20", "Animal Categories", "动物类别", 5.3, 1.5, BLUE)
    add_metric_card(slide, "0.978", "Final mAP50", "最终mAP50", 7.8, 1.5, GREEN)
    add_metric_card(slide, "90.94", "Validation Score", "验证评分", 10.3, 1.5, ORANGE)

    # Output format
    add_text_box(slide, "Output Contract / 输出格式",
                 5.3, 3.4, 4.0, 0.35, font_size=20, color=INK, bold=True)
    add_text_box(slide, '{"cat": 2, "duck": 1, "deer": 1}',
                 5.3, 3.85, 6.0, 0.45, font_size=22, color=BLUE, bold=True)

    add_bullet_list(slide, [
        ("One image in → one clean dictionary out", "单张图片输入 → 干净的字典输出"),
        ("Only supported animal categories retained", "仅保留支持的动物类别"),
        ("Counts after detection filtering & duplicate suppression", "检测过滤和重复抑制后计数"),
    ], 5.3, 4.6, 6.5, font_size=14)

    add_footer(slide, 1)
    return slide


# ============================================================
# SLIDE 2: TASK REQUIREMENTS / 任务要求
# ============================================================
def slide02_task(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, PAPER)
    add_kicker(slide, "TASK REQUIREMENTS", "任务要求")
    add_title(slide, "Structured Detection-to-Dictionary Problem",
              "结构化的检测到字典转换问题",
              "The detector must satisfy both visual recognition and strict JSON formatting.",
              "检测器需同时满足视觉识别和严格的JSON格式要求。")

    # Pipeline
    steps = [
        (1, "Detect", "检测", "Find & localize every animal", "发现并定位每个动物", BLUE),
        (2, "Filter", "过滤", "Keep only 20 allowed categories", "仅保留20个允许类别", CYAN),
        (3, "Suppress", "抑制", "Remove duplicate predictions", "移除重复预测框", ORANGE),
        (4, "Count", "计数", "Group remaining boxes by class", "按类别分组统计", GREEN),
        (5, "Export", "导出", "Write batch JSON dictionaries", "写入批量JSON字典", RED),
    ]
    for i, (num, en_l, zh_l, en_d, zh_d, c) in enumerate(steps):
        add_pipeline_step(slide, num, en_l, zh_l, en_d, zh_d, 0.6 + i * 2.45, 2.6, c)

    # Scoring focus
    add_text_box(slide, "Scoring Focus / 评分重点",
                 0.6, 4.5, 3.5, 0.35, font_size=20, color=INK, bold=True)
    add_bullet_list(slide, [
        ("Category recall: does every true species appear?", "类别召回率：是否每个真实物种都被检测到？"),
        ("Counting accuracy: is each class count correct?", "计数准确性：每个类别的计数是否正确？"),
        ("False positives: each extra wrong category costs 5 points", "误报惩罚：每个额外的错误类别扣5分"),
    ], 0.6, 5.0, 11.0, font_size=14)

    # Evaluation metrics
    add_text_box(slide, "Evaluation Metrics / 评估指标",
                 0.6, 6.1, 3.5, 0.35, font_size=16, color=INK, bold=True)
    add_text_box(slide,
                 "Format Validity / 格式有效性  |  Category Precision / 类别精确度  |  "
                 "Category Recall / 类别召回率  |  Counting Accuracy / 计数准确度  |  "
                 "Overall Response Quality / 整体响应质量",
                 0.6, 6.45, 12.0, 0.35, font_size=10, color=MUTED)

    add_footer(slide, 2)
    return slide


# ============================================================
# SLIDE 3: DATASET / 数据集构建
# ============================================================
def slide03_dataset(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, PAPER)
    add_kicker(slide, "DATASET CONSTRUCTION", "数据集构建")
    add_title(slide, "Training data expanded in stages without reducing class vocabulary.",
              "分阶段扩展训练数据，保留完整20类词汇。",
              "Every fine-tuning stage preserved the 20-class output head.",
              "每个微调阶段均保留20类输出头。")

    # Data stages table
    stages = [
        ("1. Merge / 合并", "xfz + yzy datasets", "Convert labels to JSON / 标签转换为JSON格式",
         BLUE),
        ("2. Train / 训练", "animal_dataset1 (20cls)", "20-class YOLO baseline / 20类YOLO基线模型",
         CYAN),
        ("3. Improve / 优化", "animal_dataset_5sp_v1",
         "Target: wolf/goose/horse/cow/monkey / 针对弱类别优化", ORANGE),
        ("4. Continue / 继续", "animal_dataset_6sp_v1",
         "Further optimization from best model / 从最佳模型继续优化", GREEN),
    ]
    for i, (label, data, desc, color) in enumerate(stages):
        top = 2.5 + i * 0.9
        add_rounded_rect(slide, 0.6, top, 11.6, 0.65, WHITE, LINE)
        # Left accent
        add_shape_rect(slide, 0.6, top, 0.06, 0.65, color)
        add_text_box(slide, label, 0.85, top + 0.15, 2.8, 0.3,
                     font_size=14, color=color, bold=True)
        add_text_box(slide, data, 3.8, top + 0.15, 3.5, 0.3,
                     font_size=14, color=INK, bold=True)
        add_text_box(slide, desc, 7.5, top + 0.15, 4.5, 0.3,
                     font_size=12, color=MUTED)

    # Label format info
    add_text_box(slide, "Label Format / 标签格式",
                 0.6, 6.3, 2.5, 0.3, font_size=16, color=INK, bold=True)
    add_text_box(slide,
                 "LabelMe-style JSON for annotation review → YOLO TXT format for training "
                 "/ LabelMe风格JSON用于标注审查 → YOLO TXT格式用于训练",
                 3.3, 6.3, 9.0, 0.5, font_size=12, color=MUTED)

    add_footer(slide, 3)
    return slide


# ============================================================
# SLIDE 4: MODEL PIPELINE / 模型流程
# ============================================================
def slide04_pipeline(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, PAPER)
    add_kicker(slide, "MODEL PIPELINE", "模型流程")
    add_title(slide, "YOLO inference coupled with task-specific post-processing.",
              "YOLO推理与任务特定的后处理相结合。",
              "The post-processing layer is critical because grading targets a dictionary, not box list.",
              "后处理层至关重要，因为评分目标是字典而非边界框列表。")

    # Inference pipeline
    inf_steps = [
        (1, "YOLO11s", "YOLO11s", "COCO pre-trained detector / COCO预训练检测器", BLUE),
        (2, "Fine-tune", "微调", "dataset1 → 5sp → 6sp stages / 三阶段微调", CYAN),
        (3, "Whitelist", "白名单", "20 supported labels only / 仅保留20个支持标签", GREEN),
        (4, "NMS+", "NMS+", "Balanced close-box suppression / 平衡式近距框抑制", ORANGE),
        (5, "JSON", "JSON", "Positive counts only / 仅正计数输出", RED),
    ]
    for i, (num, en_l, zh_l, en_d, c) in enumerate(inf_steps):
        add_pipeline_step(slide, num, en_l, zh_l, en_d, "", 0.6 + i * 2.45, 2.5, c)

    # Training setup
    add_text_box(slide, "Training Configuration / 训练配置",
                 0.6, 4.3, 3.5, 0.35, font_size=18, color=INK, bold=True)
    add_bullet_list(slide, [
        ("CUDA environment: NVIDIA RTX 3070 Ti Laptop GPU", "CUDA环境：NVIDIA RTX 3070 Ti 笔记本GPU"),
        ("Hyperparameters: imgsz=640, batch=8, freeze=10, patience=15",
         "超参数：图像尺寸640，批次大小8，冻结前10层，早停耐心值15"),
        ("Augmentations: HSV jitter, translation, scale, flip, mosaic, light mixup",
         "数据增强：HSV抖动、平移、缩放、翻转、马赛克、轻量mixup"),
        ("Optimizer: AdamW, lr0=0.0005, weight_decay=0.0005, warmup_epochs=3",
         "优化器：AdamW，初始学习率0.0005，权重衰减0.0005，预热3轮"),
    ], 0.6, 4.75, 11.5, font_size=13)

    add_footer(slide, 4)
    return slide


# ============================================================
# SLIDE 5: TRAINING RESULTS / 训练结果
# ============================================================
def slide05_training(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, PAPER)
    add_kicker(slide, "TRAINING RESULTS", "训练结果")
    add_title(slide, "Final fine-tuning converged quickly; epoch 2 selected as best.",
              "最终微调快速收敛，第2轮被选为最佳模型。",
              "Validation metrics remained high while later epochs showed no sustained improvement.",
              "验证指标保持高位，后期轮次无持续提升。")

    # Training curves image placeholder
    curves_path = PROJECT_ROOT / "runs" / "detect" / "runs" / "detect" / "animal_dataset_6sp_20cls_finetune" / "results.png"
    if curves_path.exists():
        slide.shapes.add_picture(str(curves_path), Inches(0.5), Inches(2.3), Inches(7.5), Inches(4.0))

    # Metrics
    add_metric_card(slide, "0.963", "Precision", "精确率", 8.6, 2.5, BLUE)
    add_metric_card(slide, "0.977", "Recall", "召回率", 10.9, 2.5, GREEN)
    add_metric_card(slide, "0.978", "mAP50", "mAP50", 8.6, 4.3, ORANGE)
    add_metric_card(slide, "0.940", "mAP50-95", "mAP50-95", 10.9, 4.3, RED)

    add_text_box(slide,
                 "Best checkpoint: runs/detect/.../animal_dataset_6sp_20cls_finetune/weights/best.pt",
                 8.6, 6.2, 4.0, 0.6, font_size=11, color=MUTED)

    # Training stages table
    add_text_box(slide, "Training Stages / 训练阶段",
                 0.5, 6.5, 3.0, 0.3, font_size=15, color=INK, bold=True)
    stages_info = [
        ("Stage 1", "animal_dataset1", "20-class baseline / 20类基线", "COCO pre-trained"),
        ("Stage 2", "animal_dataset_5sp_v1", "5 weak classes improved / 5弱类优化", "Stage 1 best.pt"),
        ("Stage 3", "animal_dataset_6sp_v1", "Continued optimization / 继续优化", "Stage 2 best.pt"),
    ]
    for i, (stage, data, purpose, base) in enumerate(stages_info):
        top = 6.8
        left = 0.5 + i * 4.1
        add_text_box(slide, f"{stage}: {data}", left, top, 3.8, 0.2,
                     font_size=10, color=BLUE, bold=True)
        add_text_box(slide, f"{purpose}  ←  {base}", left, top + 0.2, 3.8, 0.2,
                     font_size=9, color=MUTED)

    add_footer(slide, 5)
    return slide


# ============================================================
# SLIDE 6: VALIDATION RESULTS / 验证结果
# ============================================================
def slide06_validation(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, PAPER)
    add_kicker(slide, "VALIDATION RESULTS", "验证结果")
    add_title(slide, "Balanced suppression produced 90.94/100 average validation score.",
              "平衡式抑制策略获得了90.94/100的平均验证评分。",
              "Remaining errors concentrated in confusing bird & canine-like classes.",
              "剩余错误集中在易混淆的禽类和犬科类别。")

    # Key metrics
    add_metric_card(slide, "90.94", "Avg Score / 均分", "/100", 0.6, 2.5, ORANGE)
    add_metric_card(slide, "9/15", "Perfect Images", "满分图片", 3.1, 2.5, GREEN)
    add_metric_card(slide, "15", "Validation Set", "验证集图片数", 5.6, 2.5, BLUE)

    # Score bars
    add_text_box(slide, "Per-Image Scores / 逐图评分",
                 0.6, 4.3, 3.0, 0.35, font_size=16, color=INK, bold=True)
    scores = [
        ("eval_001384", 100), ("eval_001388", 100), ("eval_001431", 100),
        ("val_4", 100), ("val_5", 100), ("val_6", 100),
        ("val_7", 100), ("val_8", 100), ("val_10", 100),
        ("eval_001391", 83.33),
        ("val_1", 75), ("val_2", 61.67), ("val_9", 61.67),
        ("eval_001402", 50), ("val_3", 50),
    ]
    for i, (name, score) in enumerate(scores):
        col = i % 2
        row = i // 2
        left = 0.6 + col * 6.2
        top = 4.8 + row * 0.32
        bar_color = GREEN if score >= 100 else (ORANGE if score >= 70 else RED)
        add_bar_chart_row(slide, name, score, 100, left, top, bar_color)

    # Error analysis
    add_text_box(slide, "Main Error Types / 主要错误类型",
                 8.4, 2.6, 4.0, 0.3, font_size=16, color=INK, bold=True)
    add_bullet_list(slide, [
        ("Missed detections (e.g., fox in val_1)", "漏检（如val_1中的狐狸）"),
        ("Class confusion: duck ↔ goose", "类别混淆：鸭 ↔ 鹅"),
        ("Class confusion: fox/wolf/dog", "类别混淆：狐狸/狼/狗"),
        ("Extra false positive categories", "额外的误报类别"),
        ("Count errors in occluded scenes", "遮挡场景中的计数错误"),
    ], 8.4, 3.1, 4.0, font_size=12)

    add_footer(slide, 6)
    return slide


# ============================================================
# SLIDE 7: VISUAL OUTPUT / 可视化输出
# ============================================================
def slide07_visual(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, PAPER)
    add_kicker(slide, "VISUAL OUTPUT", "可视化输出")
    add_title(slide, "Validation images annotated with boxes & confidence scores.",
              "验证集图片标注了边界框和置信度分数。",
              "Same post-processing logic used for both displayed boxes and JSON counts.",
              "显示框和JSON计数使用相同的后处理逻辑。")

    # Annotated image
    anno_path = PROJECT_ROOT / "outputs" / "val_set_annotated_6sp_20cls_finetuned_suppressed_balanced" / "annotated" / "eval_001384.png"
    if anno_path.exists():
        slide.shapes.add_picture(str(anno_path), Inches(0.6), Inches(2.3), Inches(5.8), Inches(4.3))

    # Post-processing rules
    add_text_box(slide, "Post-Processing Rules / 后处理规则",
                 7.0, 2.5, 5.0, 0.35, font_size=18, color=INK, bold=True)
    add_bullet_list(slide, [
        ("Same-class: merge when IoU high or centers close with similar size",
         "同类别：IoU高或中心距近且尺寸相似时合并"),
        ("Cross-class: merge only when likely covering same object (stricter thresholds)",
         "跨类别：仅在很可能覆盖同一目标时合并（更严格阈值）"),
        ("Keep highest-confidence detection among duplicates",
         "在重复检测中保留置信度最高的"),
        ("Balanced mode: reduce duplicates without aggressive deletion",
         "平衡模式：减少重复但不激进删除邻近动物"),
    ], 7.0, 3.0, 5.5, font_size=12)

    add_text_box(slide, "Output Directory / 输出目录",
                 7.0, 5.5, 3.0, 0.3, font_size=14, color=BLUE, bold=True)
    add_text_box(slide,
                 "outputs/val_set_annotated_6sp_20cls_finetuned_suppressed_balanced/\n"
                 "  annotated/  ← images with boxes / 带框图片\n"
                 "  detections_with_boxes.json  ← full detection metadata / 完整检测元数据",
                 7.0, 5.85, 5.5, 0.8, font_size=10, color=MUTED)

    add_footer(slide, 7)
    return slide


# ============================================================
# SLIDE 8: ANALYSIS & CHALLENGES / 分析与挑战
# ============================================================
def slide08_analysis(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, PAPER)
    add_kicker(slide, "ANALYSIS & CHALLENGES", "分析与挑战")
    add_title(slide, "Deep dive into model behavior and remaining issues.",
              "深入分析模型行为和遗留问题。")

    # Strengths
    add_rounded_rect(slide, 0.5, 2.2, 3.8, 2.5, WHITE, LINE)
    add_shape_rect(slide, 0.5, 2.2, 3.8, 0.06, GREEN)
    add_text_box(slide, "Strengths / 优势", 0.7, 2.4, 3.4, 0.3,
                 font_size=18, color=GREEN, bold=True)
    add_bullet_list(slide, [
        ("Large distinctive animals (zebra, giraffe, horse, lion, tiger, bear) detected well",
         "大型特征明显动物检测良好"),
        ("High mAP50 (0.978) and mAP50-95 (0.940)",
         "高mAP50和mAP50-95指标"),
        ("Three-stage fine-tuning improved weak classes",
         "三阶段微调改善了弱类别"),
        ("Clean JSON output format, no NL explanations",
         "干净的JSON输出格式"),
    ], 0.7, 2.85, 3.4, font_size=11)

    # Weaknesses
    add_rounded_rect(slide, 4.7, 2.2, 3.8, 2.5, WHITE, LINE)
    add_shape_rect(slide, 4.7, 2.2, 3.8, 0.06, RED)
    add_text_box(slide, "Limitations / 局限", 4.9, 2.4, 3.4, 0.3,
                 font_size=18, color=RED, bold=True)
    add_bullet_list(slide, [
        ("Duck ↔ goose confusion (similar body shape & color)",
         "鸭鹅混淆（体型颜色相似）"),
        ("Fox/wolf/dog confusion (similar texture & pose)",
         "狐狸/狼/狗混淆（纹理姿态相似）"),
        ("Occlusion & tightly grouped animals",
         "遮挡和密集排列动物"),
        ("Synthetic scene generalization",
         "合成场景泛化能力"),
    ], 4.9, 2.85, 3.4, font_size=11)

    # Solutions
    add_rounded_rect(slide, 8.9, 2.2, 3.9, 2.5, WHITE, LINE)
    add_shape_rect(slide, 8.9, 2.2, 3.9, 0.06, BLUE)
    add_text_box(slide, "Solutions / 解决方案", 9.1, 2.4, 3.5, 0.3,
                 font_size=18, color=BLUE, bold=True)
    add_bullet_list(slide, [
        ("Custom post-processing for duplicate suppression",
         "自定义后处理抑制重复框"),
        ("Targeted fine-tuning on weak-class datasets",
         "针对弱类数据集定向微调"),
        ("Balanced thresholds to avoid aggressive filtering",
         "平衡阈值避免激进过滤"),
        ("Maintained 20-class head throughout all stages",
         "全程保持20类输出头"),
    ], 9.1, 2.85, 3.5, font_size=11)

    # Key insight box
    add_rounded_rect(slide, 0.5, 5.0, 12.3, 1.8, LIGHT_BLUE_BG, BLUE)
    add_text_box(slide, "Key Findings / 关键发现",
                 0.7, 5.15, 4.0, 0.3, font_size=16, color=BLUE, bold=True)
    add_bullet_list(slide, [
        ("Maintaining 20-class output during fine-tuning prevented vocabulary loss while improving weak categories",
         "微调时保持20类输出避免了词汇丢失同时改善了弱类别"),
        ("Balanced post-processing (IoU=0.40 for same-class, 0.62 for cross-class) reduced duplicates without hurting recall",
         "平衡后处理（同类别IoU=0.40，跨类别IoU=0.62）减少重复但未损害召回率"),
        ("Scoring rule strongly penalizes false positive categories (5 pts each) → conservative threshold tuning needed",
         "评分规则对误报类别惩罚严厉（每类5分）→ 需要保守的阈值调优"),
    ], 0.7, 5.55, 11.8, font_size=12)

    add_footer(slide, 8)
    return slide


# ============================================================
# SLIDE 9: CONTRIBUTIONS / 贡献
# ============================================================
def slide09_contributions(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, PAPER)
    add_kicker(slide, "TEAM CONTRIBUTIONS", "团队贡献")
    add_title(slide, "Each member's role and contribution to the project.",
              "每位成员在项目中的角色与贡献。")

    # Team member cards
    members = [
        ("Member 1 / 成员1", "待填写 / TBD",
         [
             "Dataset collection & merging / 数据集收集与合并",
             "LabelMe JSON annotation workflow / LabelMe JSON标注流程",
             "Validation data organization / 验证数据组织",
             "Data preprocessing scripts / 数据预处理脚本",
         ], BLUE),
        ("Member 2 / 成员2", "待填写 / TBD",
         [
             "YOLO model training & fine-tuning / YOLO模型训练与微调",
             "CUDA environment setup & optimization / CUDA环境配置与优化",
             "Model metric tracking & selection / 模型指标跟踪与选型",
             "Hyperparameter tuning / 超参数调优",
         ], GREEN),
        ("Member 3 / 成员3", "待填写 / TBD",
         [
             "Evaluation & inference scripts / 评估与推理脚本",
             "Duplicate-box post-processing algorithm / 重复框后处理算法",
             "Report, PPT & documentation / 报告、PPT与文档",
             "Error analysis & validation scoring / 错误分析与验证评分",
         ], ORANGE),
    ]

    for i, (name, id_str, tasks, color) in enumerate(members):
        left = 0.5 + i * 4.2
        add_rounded_rect(slide, left, 2.3, 3.9, 4.3, WHITE, LINE)
        add_shape_rect(slide, left, 2.3, 3.9, 0.06, color)
        add_text_box(slide, name, left + 0.2, 2.5, 3.5, 0.35,
                     font_size=18, color=color, bold=True)
        add_text_box(slide, f"ID: {id_str}", left + 0.2, 2.85, 3.5, 0.25,
                     font_size=12, color=MUTED)
        add_shape_rect(slide, left + 0.2, 3.2, 3.5, 0.01, LINE)
        add_bullet_list(slide, tasks, left + 0.2, 3.35, 3.5, font_size=12)

    add_footer(slide, 9)
    return slide


# ============================================================
# SLIDE 10: CONCLUSION / 结论
# ============================================================
def slide10_conclusion(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, DARK_BG)

    # Title
    add_text_box(slide, "Conclusion / 结论",
                 0.6, 0.8, 6.0, 0.7, font_size=38, color=WHITE, bold=True)

    add_text_box(slide,
                 "A reproducible YOLO-based animal counting system covering data preparation, "
                 "fine-tuning, validation scoring, JSON export, and visual annotation.",
                 0.7, 1.7, 10.0, 0.65, font_size=18,
                 color=RGBColor(0xDC, 0xE8, 0xF8))
    add_text_box(slide,
                 "可复现的基于YOLO的动物计数系统，涵盖数据准备、微调、验证评分、JSON导出和可视化标注。",
                 0.7, 2.15, 10.0, 0.4, font_size=14,
                 color=RGBColor(0xBB, 0xD2, 0xF0))

    # Three cards
    cards = [
        ("Achievements / 成果", GREEN,
         [
             "High mAP50 of 0.978 across 20 animal classes",
             "20个动物类别mAP50达0.978",
             "Validation score of 90.94/100 on 15-image set",
             "15张验证集评分90.94/100",
             "Clean dictionary-format JSON output",
             "干净的字典格式JSON输出",
             "Complete code + report + PPT deliverables",
             "完整的代码+报告+PPT交付物",
         ]),
        ("Limitations / 局限", ORANGE,
         [
             "Duck/goose and fox/wolf/dog confusions persist",
             "鸭/鹅和狐狸/狼/狗的混淆仍然存在",
             "Occlusion & tight grouping cause count errors",
             "遮挡和密集排列导致计数错误",
             "Synthetic-to-real domain gap possible",
             "合成到真实场景可能存在域差距",
         ]),
        ("Next Steps / 后续工作", BLUE,
         [
             "Add manually verified hard examples for weak classes",
             "增加弱类别的人工验证难例",
             "Continue targeted fine-tuning with stricter review",
             "以更严格审查继续定向微调",
             "Explore test-time augmentation for robustness",
             "探索测试时增强提升鲁棒性",
             "Improve duck/goose discrimination features",
             "改善鸭/鹅区分特征",
         ]),
    ]

    for i, (title, color, items) in enumerate(cards):
        left = 0.5 + i * 4.2
        add_rounded_rect(slide, left, 2.8, 3.9, 3.8,
                         RGBColor(0x14, 0x2A, 0x46), RGBColor(0x35, 0x5B, 0x88))
        add_text_box(slide, title, left + 0.2, 3.0, 3.5, 0.35,
                     font_size=18, color=color, bold=True)
        txBox = slide.shapes.add_textbox(Inches(left + 0.2), Inches(3.5), Inches(3.5), Inches(2.8))
        tf = txBox.text_frame
        tf.word_wrap = True
        for j, item in enumerate(items):
            if j == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            if j % 2 == 0:
                p.text = f"▸ {item}"
                p.font.size = Pt(13)
                p.font.color.rgb = WHITE
                p.font.name = "Arial"
            else:
                p.text = f"    {item}"
                p.font.size = Pt(11)
                p.font.color.rgb = RGBColor(0xDC, 0xE8, 0xF8)
                p.font.name = "Microsoft YaHei"
            p.space_after = Pt(3)

    # Deliverables footer
    add_text_box(slide,
                 "Final Deliverables / 最终交付物: Report / 报告  |  PPT  |  Source Code / 源代码  |  "
                 "Evaluation JSON / 评估JSON  |  Annotated Images / 标注图片",
                 0.7, 6.9, 11.5, 0.35, font_size=13,
                 color=RGBColor(0xBB, 0xD2, 0xF0))

    add_footer(slide, 10)
    return slide


# ============================================================
# MAIN
# ============================================================
def main():
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide01_cover(prs)
    slide02_task(prs)
    slide03_dataset(prs)
    slide04_pipeline(prs)
    slide05_training(prs)
    slide06_validation(prs)
    slide07_visual(prs)
    slide08_analysis(prs)
    slide09_contributions(prs)
    slide10_conclusion(prs)

    prs.save(str(PPT_PATH))
    print(f"PPT saved to: {PPT_PATH}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
